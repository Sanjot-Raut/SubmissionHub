import fitz  # PyMuPDF
from sentence_transformers import SentenceTransformer, util
from pptx import Presentation
from docx import Document
import mysql.connector
from mysql.connector import Error
import json

# Function to extract text from a PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text

# Function to extract text from a PowerPoint file
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Function to extract text from a Word document
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text

# Function to determine file type and extract text accordingly
def extract_text(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    else:
        raise ValueError("Unsupported file format. Please provide a PDF, PPTX, or DOCX file.")

# Load pre-trained model for semantic similarity
model = SentenceTransformer('all-MiniLM-L6-v2')

# Function to vectorize text
def vectorize_text(text):
    return model.encode(text, convert_to_tensor=True)

# Function to compare two vectorized texts using cosine similarity
def compare_texts(vector1, vector2):
    similarity_score = util.pytorch_cos_sim(vector1, vector2).item()
    return similarity_score

# Function to get existing projects from the database
def get_existing_projects(cursor):
    cursor.execute("SELECT id, vectorized_text FROM projects WHERE vectorized_text IS NOT NULL")
    return cursor.fetchall()

# Function to save new project to the database
def save_new_project(connection,cursor, title, description, student_id, assignment_id,  file_path, vectorized_text):
    cursor.execute(
        "INSERT INTO projects (title, description, student_id, assignment_id, file_path, vectorized_text) "
        "VALUES (%s, %s, %s, %s, %s, %s)",
        (title, description, student_id, assignment_id,  file_path, vectorized_text)
    )
        # Commit the transaction to ensure the data is saved
    connection.commit()
    
    # Return the ID of the inserted row
    return cursor.lastrowid


def save_comparison_result(connection,cursor, project_id, status,similar_id):
    similar_to_id=json.dumps(similar_id)
    cursor.execute(
        "INSERT INTO project_comparisons (project_id,status,similar_to_id) "
        "VALUES (%s, %s, %s)",
        (project_id,status,similar_to_id)
    )
    connection.commit()
# Main function to process new submission
def process_new_submission(title, description, student_id, assignment_id,  file_path):
   
    try:
        connection = mysql.connector.connect(
            host='localhost',
            database='submission_hub',
            user='root',
            password=''
        )
        cursor = connection.cursor()
        # Extract and vectorize new project text
        new_text = extract_text(file_path)
        new_vector = vectorize_text(new_text).tolist()  # Convert tensor to list for storage

        # Compare with existing projects
        existing_projects = get_existing_projects(cursor)
        similarities = [(proj_id, compare_texts(new_vector, [float(x) for x in existing_vector.decode('utf-8')[1:-1].split(',')]))
                        for proj_id, existing_vector in existing_projects if existing_vector]

        # Determine if the project is unique or copied
        total_threshold = 0.9
        partial_threshold = 0.5
        total_copies = [proj_id for proj_id, score in similarities if score >= total_threshold]
        partial_copies = [proj_id for proj_id, score in similarities if partial_threshold <= score < total_threshold]
        vectorized_text_str = str(new_vector) 

        if total_copies:
            project_id=save_new_project(connection,cursor, title, description, student_id, assignment_id,  file_path, vectorized_text_str)
            save_comparison_result(connection,cursor, project_id, "totally copied",total_copies)
            return "totally copied", total_copies
        elif partial_copies:
            project_id=save_new_project(connection,cursor, title, description, student_id, assignment_id,  file_path, vectorized_text_str)
            save_comparison_result(connection,cursor, project_id,"partially copied",partial_copies)
            return "partially copied", partial_copies
        else:
            project_id=save_new_project(connection,cursor, title, description, student_id, assignment_id,  file_path, vectorized_text_str)
            save_comparison_result(connection,cursor, project_id, "unique",[])
            connection.commit()
            return "unique", []
         
    except Error as e:
        print(f"Error: {e}")
        return None, None
    finally:
        if connection is not None and connection.is_connected():
            cursor.close()
            connection.close()

def project_submission_process(title,description,student_id,assignment_id,file_path):

    status, copies = process_new_submission(title, description, student_id, assignment_id, file_path)
    
    if status == "unique":
        print("The project is unique and has been added to the database.")
    else:
        print(f"The project is {status}. Similar projects found:")
        for copy in copies:
            print(f"Project ID: {copy}")
